import io
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import plotly.graph_objects as go
from datetime import datetime
from typing import List
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import contextlib


class PowerPointExporter:
    def __init__(self):
        # colors and sizing for the powerpoint presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

        self.colors = {
            "primary": RGBColor(51, 80, 184),  # 3350B8
            "secondary": RGBColor(0, 0, 0),
            "bottom_box": RGBColor(15, 30, 50),  # 0F1E32
            "accent": RGBColor(255, 127, 0),
            "text": RGBColor(255, 255, 255),
        }

        self.font_name = "Arial"

    def add_title_slide(self, title: str, subtitle: str = None):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            self.prs.slide_width,
            self.prs.slide_height,
        )
        fill = bg.fill
        fill.gradient()

        fill.gradient_stops[0].color.rgb = self.colors["secondary"]  # starts with blue on the top
        fill.gradient_stops[0].position = 0.0

        fill.gradient_stops[1].color.rgb = self.colors["primary"]  # fades to black on the bottom
        fill.gradient_stops[1].position = 0.5
        bg.line.fill.background()

        # title textbox
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.font_name
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.colors["text"]

        # subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
            tf_sub = sub_box.text_frame
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.name = self.font_name
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = self.colors["text"]

    def add_image_slide(self, image_path: str, title: str = None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height
        box_h = slide_h / 3

        # background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.fill.background()

        # bottom box on the background
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            slide_h - box_h,
            slide_w,
            box_h,
        )
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = self.colors["bottom_box"]
        bottom.line.fill.background()

        # constants
        title_space = Inches(0.8)
        img_left = Inches(0.5)
        max_width = slide_w - Inches(1)

        # add image
        pic = slide.shapes.add_picture(
            image_path,
            img_left,
            title_space,
            width=max_width,
        )
        scale = 0.92
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)

        # re-center horizontally after scaling
        pic.left = int((slide_w - pic.width) / 2 - Inches(0.15))
        pic.top = int(title_space)

        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), slide_w - Inches(1), Inches(0.8))
            p = title_box.text_frame.paragraphs[0]
            p.text = title
            p.font.name = self.font_name
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.colors["bottom_box"]

    def save(self):
        import io

        pptx_bytes = io.BytesIO()
        self.prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes


def export_figures_to_temp_files(figures: List[go.Figure]) -> List[str]:
    # exports figures to a temporary file outside of the main button callback
    # done because direct exporting creates a conflict with streamlit and kaleido
    temp_files = []

    for i, fig in enumerate(figures):
        try:
            temp_fd, temp_path = tempfile.mkstemp(suffix=".png", prefix=f"chart_{i}_")
            os.close(temp_fd)

            fig.write_image(temp_path, format="png", width=1400, height=788, scale=2)

            temp_files.append(temp_path)
        except Exception as e:  # noqa: BLE001, PERF203
            st.error(f"Failed to export chart {i + 1}: {e!s}")
            # placeholder file path
            temp_files.append(None)

    return temp_files


def create_ppt_from_images(image_paths: List[str], page_name: str, layout: str) -> io.BytesIO:  # noqa: D103
    # creates a powerpoint from the pre generated image files
    exporter = PowerPointExporter()
    exporter.add_title_slide(title=f"{page_name}", subtitle=f"Generated on {datetime.now():%B %d, %Y}")  # noqa: DTZ005
    for i, img_path in enumerate(image_paths):
        if img_path and os.path.exists(img_path):  # noqa: PTH110
            exporter.add_image_slide(image_path=img_path, title=f"{page_name} - Chart {i + 1}")

    return exporter.save()


def cleanup_temp_files(file_paths: List[str]):
    # delete the temporary intermediary files
    for path in file_paths:
        if path and os.path.exists(path):  # noqa: PTH110
            with contextlib.suppress(BaseException):
                os.unlink(path)  # noqa: PTH108


def create_export_button(figures_col1, figures_col2, page_name="Performance"):
    # creates the exporting button to the page with the logic of generating
    # the images before the powerpoint generation button
    with st.sidebar:
        st.markdown("---")
        st.subheader("Export to PowerPoint")

        all_figures = figures_col1 + figures_col2
        total = len(all_figures)
        st.caption(f"{total} charts ready")

        # Session state keys
        images_key = f"images_{page_name}"
        ppt_key = f"ppt_{page_name}"

        # Initialize session state
        if images_key not in st.session_state:
            st.session_state[images_key] = None
        if ppt_key not in st.session_state:
            st.session_state[ppt_key] = None

        # pre-generating the images of all the charts on the page
        if st.button("Prepare Charts", use_container_width=True, type="primary"):
            with st.spinner(f"Exporting {total} charts to images..."):
                try:
                    temp_files = export_figures_to_temp_files(all_figures)
                    st.session_state[images_key] = temp_files

                    valid_count = sum(1 for f in temp_files if f)
                    st.success(f"Prepared {valid_count}/{total} charts")

                except Exception as e:
                    st.error(f"Failed: {str(e)}")

        # creates powerpoint, if images of the charts are done
        if st.session_state[images_key]:
            if st.button("Generate PowerPoint", use_container_width=True, type="secondary"):
                with st.spinner("Creating PowerPoint..."):
                    try:
                        ppt_bytes = create_ppt_from_images(st.session_state[images_key], page_name, "single")
                        st.session_state[ppt_key] = ppt_bytes

                        cleanup_temp_files(st.session_state[images_key])
                        st.session_state[images_key] = None

                        st.success("PowerPoint is ready")

                    except Exception as e:
                        st.error(f"Failed: {str(e)}")

        if st.session_state[ppt_key]:
            st.success("Ready to download")

            filename = f"{page_name}_Charts_{datetime.now():%Y%m%d_%H%M%S}.pptx"

            st.download_button(
                "Download PowerPoint",
                data=st.session_state[ppt_key],
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )
            if st.button("Export Again", use_container_width=True):
                st.session_state[ppt_key] = None
                st.rerun()
